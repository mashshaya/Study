from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import copy

FINAL = Path("/Users/maria/Desktop/Code/HSE/COURSEBOOK/final")
FIGS  = FINAL / "results_models" / "figures_final"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]
W = prs.slide_width
H = prs.slide_height

# ── Palette ──────────────────────────────────────────────────────────────────
FONT         = "Calibri Light"
C_TITLE      = RGBColor(0x14, 0x4F, 0x7B)   # dark readable blue for titles
C_BOX_MAIN   = RGBColor(0x1E, 0x77, 0xB4)   # medium blue — main boxes
C_BOX_ALT    = RGBColor(0x45, 0x9A, 0xCC)   # lighter blue — secondary boxes
C_BOX_GREEN  = RGBColor(0x1D, 0x6A, 0x4A)   # green — conclusions / good news
C_BOX_RED    = RGBColor(0xC0, 0x39, 0x2B)   # red — warnings / contrast facts
C_TBL_HEAD   = RGBColor(0x1E, 0x77, 0xB4)   # table header fill
C_TBL_ROW    = RGBColor(0xEB, 0xF5, 0xFB)   # table alternating row
C_TEXT       = RGBColor(0x1A, 0x1A, 0x2E)   # body text
C_ARROW      = RGBColor(0x6A, 0x6A, 0x6A)   # arrows in schemas
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def _set_font(p, size, bold=False, color=None, name=FONT):
    p.font.name  = name
    p.font.size  = Pt(size)
    p.font.bold  = bold
    if color:
        p.font.color.rgb = color

def add_title(slide, text, top=Inches(0.22), font_size=30):
    tb = slide.shapes.add_textbox(Inches(0.4), top, W - Inches(0.8), Inches(0.72))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _set_font(p, font_size, bold=True, color=C_TITLE)
    # thin accent line under title — via a narrow rectangle
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                  Inches(0.4), top + Inches(0.68),
                                  W - Inches(0.8), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = C_BOX_ALT
    line.line.fill.background()
    return tb

def add_body(slide, lines, left=Inches(0.4), top=Inches(1.1),
             width=None, height=None, font_size=17, bold_first=False):
    if width  is None: width  = W - Inches(0.8)
    if height is None: height = H - top - Inches(0.15)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        stripped = line.lstrip()
        indent   = len(line) - len(stripped)
        bullet   = stripped.startswith("- ") or stripped.startswith("• ")
        text     = stripped[2:] if bullet else stripped
        p.text   = text
        _set_font(p, font_size, bold=(bold_first and first), color=C_TEXT)
        if bullet:
            p.level  = 1 if indent >= 4 else 0
            run      = p.runs[0] if p.runs else p.add_run()
            run.text = ("    " if indent >= 4 else "") + "• " + text
            p.text   = ""
        if not stripped:
            p.space_after = Pt(3)
    return tb

def add_image(slide, img_path, left, top, width, height=None):
    if not Path(img_path).exists():
        return
    if height:
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(img_path), left, top, width)

def add_table(slide, headers, rows, left, top, width, height,
              font_size=15, header_color=None):
    if header_color is None:
        header_color = C_TBL_HEAD
    cols = len(headers)
    tbl  = slide.shapes.add_table(len(rows)+1, cols, left, top, width, height).table
    col_w = width // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        _set_font(p, font_size, bold=True, color=C_WHITE)
        p.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            _set_font(p, font_size, color=C_TEXT)
            p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C_TBL_ROW
    return tbl

def add_box(slide, text, left, top, width, height,
            fill=None, text_color=C_WHITE, font_size=14, bold=True):
    if fill is None:
        fill = C_BOX_MAIN
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        _set_font(p, font_size, bold=bold, color=text_color)
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_down_arrow(slide, left, top, width, h=Inches(0.28)):
    tb = slide.shapes.add_textbox(left, top, width, h)
    p  = tb.text_frame.paragraphs[0]
    p.text = "↓"; p.alignment = PP_ALIGN.CENTER
    _set_font(p, 18, bold=True, color=C_ARROW)

def add_arrow(slide, left, top, height=Inches(0.5)):
    tb = slide.shapes.add_textbox(left, top, Inches(0.45), height)
    p  = tb.text_frame.paragraphs[0]
    p.text = "→"; p.alignment = PP_ALIGN.CENTER
    _set_font(p, 24, bold=True, color=C_ARROW)

def add_sublabel(slide, text, left, top, width=Inches(2.1), font_size=12):
    tb = slide.shapes.add_textbox(left, top, width, Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        _set_font(p, font_size, color=RGBColor(0x44, 0x44, 0x44))
        p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — Титульный
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()

# Thin top accent bar
bar = sl.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                           Inches(0), Inches(0), W, Inches(0.18))
bar.fill.solid(); bar.fill.fore_color.rgb = C_BOX_MAIN
bar.line.fill.background()

tb0 = sl.shapes.add_textbox(Inches(0.5), Inches(1.1), W - Inches(1), Inches(0.45))
p0  = tb0.text_frame.paragraphs[0]
p0.text = "НИУ ВШЭ · Факультет экономических наук · Образовательная программа «Экономика»"
_set_font(p0, 14, color=RGBColor(0x77, 0x77, 0x77))
p0.alignment = PP_ALIGN.CENTER

tb = sl.shapes.add_textbox(Inches(0.5), Inches(1.75), W - Inches(1), Inches(1.9))
tf = tb.text_frame; tf.word_wrap = True
p  = tf.paragraphs[0]
p.text = ("Исследование устойчивости моделей волатильности\n"
          "и ценообразования опционов\nв различных рыночных режимах")
_set_font(p, 31, bold=True, color=C_TITLE)
p.alignment = PP_ALIGN.CENTER

# Divider line
div = sl.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                           Inches(2.5), Inches(3.85), Inches(8.33), Inches(0.04))
div.fill.solid(); div.fill.fore_color.rgb = C_BOX_ALT
div.line.fill.background()

for text, sz, top_in in [
    ("Ковалёва Мария Сергеевна, 3 курс, группа БЭК2310", 17, 4.05),
    ("Научный руководитель: Сизых Наталья Васильевна",   16, 4.5),
    ("Москва, 2026",                                      15, 4.95),
]:
    tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(top_in), W - Inches(1), Inches(0.48))
    p2  = tb2.text_frame.paragraphs[0]
    p2.text = text
    _set_font(p2, sz, color=RGBColor(0x33, 0x33, 0x33))
    p2.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — Контекст и проблема
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Контекст и проблема")

for text, bx, col in [
    ("Рынок\n\nОпционы торгуются на бирже.\nУ каждого контракта есть\nнаблюдаемая рыночная цена.",
     Inches(0.4), C_BOX_MAIN),
    ("Модели\n\nДля оценки «справедливой» цены используют формульные модели:\nBlack-76, CRR и др.",
     Inches(4.55), C_BOX_ALT),
    ("Проблема\n\nМодели учитывают ограниченный набор факторов, а рыночная цена агрегирует больше информации. Поэтому ошибка зависит от режима и характеристик опциона.",
     Inches(8.7), C_BOX_MAIN),
]:
    add_box(sl, text, left=bx, top=Inches(1.15), width=Inches(3.9), height=Inches(2.7),
            fill=col, font_size=15)

add_body(sl, ["→ Насколько точны популярные модели на реальных российских данных?"],
         top=Inches(4.1), font_size=18)

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — Гипотеза и цель
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Гипотеза и цель исследования")

add_box(sl, "Гипотеза", left=Inches(0.4), top=Inches(1.1), width=Inches(12.5), height=Inches(0.48),
        fill=C_BOX_MAIN, font_size=15)
add_body(sl, [
    "Ошибка моделей зависит от рыночного режима и возрастает в stress regime",
], top=Inches(1.68), font_size=21)

add_box(sl, "Цель", left=Inches(0.4), top=Inches(2.75), width=Inches(12.5), height=Inches(0.48),
        fill=C_BOX_MAIN, font_size=15)
add_body(sl, [
    "Исследовать устойчивость моделей волатильности и ценообразования опционов",
    "в различных рыночных режимах на основе сравнения модельных и рыночных цен",
], top=Inches(3.33), font_size=18)

add_body(sl, [
    "→ Проверяем не «истинную цену», а устойчивость ошибки: как и где она меняется",
], top=Inches(4.65), font_size=16)

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — Методология
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Методология исследования")

# LEFT: vertical flow schema — compact boxes
SCHEMA_X = Inches(0.4)
SCHEMA_W = Inches(5.5)
BOX_H    = Inches(0.76)
ARR_H    = Inches(0.28)
Y        = Inches(1.1)

flow = [
    ("MOEX data",                                           C_BOX_MAIN),
    ("Оценка волатильности\nHV_21d  ·  HV_63d  ·  GARCH", C_BOX_ALT),
    ("Модель ценообразования\nCRR American  ·  Black-76",  C_BOX_ALT),
    ("Ошибка:  model price − market price",                C_BOX_GREEN),
    ("Анализ по vol_regime · DTE · moneyness",             C_BOX_GREEN),
]

for i, (label, color) in enumerate(flow):
    add_box(sl, label, left=SCHEMA_X, top=Y, width=SCHEMA_W, height=BOX_H,
            fill=color, font_size=14)
    Y += BOX_H
    if i < len(flow) - 1:
        add_down_arrow(sl, SCHEMA_X, Y, SCHEMA_W)
        Y += ARR_H

# RIGHT: models + metrics
RX = Inches(6.2)
RW = Inches(6.8)

add_body(sl, [
    "- CRR American — основное сравнение",
    "- Black-76 — европейский benchmark",
], left=RX, top=Inches(1.1), width=RW, font_size=16)

add_box(sl, "Метрики качества", left=RX, top=Inches(2.9), width=RW, height=Inches(0.42),
        fill=C_BOX_MAIN, font_size=14)
add_body(sl, ["MAE  ·  RMSE  ·  Mean error"],
         left=RX, top=Inches(3.4), width=RW, font_size=16)

add_box(sl, "Авторская метрика  TVNAE", left=RX, top=Inches(4.25), width=RW, height=Inches(0.42),
        fill=C_BOX_ALT, font_size=14)

tb_f = sl.shapes.add_textbox(RX, Inches(4.75), RW, Inches(0.55))
p_f  = tb_f.text_frame.paragraphs[0]
p_f.text = "TVNAEᵢ = |P̂ᵢ − Pᵢ| / TVᵢ"
_set_font(p_f, 18, bold=True, color=C_TITLE)
p_f.alignment = PP_ALIGN.CENTER

add_body(sl, [
    "Ошибка нормируется на временну́ю стоимость —",
    "ту часть цены, которую модель должна объяснять",
], left=RX, top=Inches(5.38), width=RW, font_size=14)

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — Модели ценообразования
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Модели ценообразования и входные параметры")

# Two model blocks side by side
add_box(sl, "CRR American",
        left=Inches(0.4), top=Inches(1.05), width=Inches(6.0), height=Inches(0.46),
        fill=C_BOX_MAIN, font_size=15)
add_body(sl, [
    "Биномиальная модель для американских опционов.",
    "Строит дерево цен, допускает досрочное исполнение.",
    "Входы: F, K, T, r, σ, option type",
], top=Inches(1.6), width=Inches(6.0), font_size=15)

add_box(sl, "Black-76",
        left=Inches(7.0), top=Inches(1.05), width=Inches(6.0), height=Inches(0.46),
        fill=C_BOX_ALT, font_size=15)
add_body(sl, [
    "Европейский benchmark для опционов на фьючерсы.",
    "Аналитическая формула, без досрочного исполнения.",
    "Входы: F, K, T, r, σ",
], top=Inches(1.6), left=Inches(7.0), width=Inches(6.0), font_size=15)

add_body(sl, [
    "σ — ключевой параметр: оценка σ варьируется между моделями.",
], top=Inches(3.0), width=Inches(12.5), font_size=14)

# Horizontal mini-schema at bottom: [HV] → [σ] → [CRR / B76] → [price]
tb_r = sl.shapes.add_textbox(Inches(0.4), Inches(3.5), W - Inches(0.8), Inches(0.55))
p_r  = tb_r.text_frame.paragraphs[0]
p_r.text = "Разные оценки σ → разные модельные цены → разные ошибки"
_set_font(p_r, 17, bold=True, color=C_BOX_RED)
p_r.alignment = PP_ALIGN.CENTER

# 4 boxes + 3 arrows horizontally
BX   = Inches(0.54);  BW = Inches(2.75); BH = Inches(0.82); BY = Inches(4.22)
ARRW = Inches(0.38)
hflow = [
    ("HV_21d · HV_63d · GARCH", C_BOX_ALT),
    ("σ",                        C_BOX_MAIN),
    ("CRR  /  Black-76",         C_BOX_ALT),
    ("model price",              C_BOX_GREEN),
]
cx = BX
for k, (lbl, col) in enumerate(hflow):
    add_box(sl, lbl, left=cx, top=BY, width=BW, height=BH, fill=col, font_size=15)
    cx += BW
    if k < len(hflow) - 1:
        tb_arr = sl.shapes.add_textbox(cx, BY, ARRW, BH)
        p_arr  = tb_arr.text_frame.paragraphs[0]
        p_arr.text = "→"; p_arr.alignment = PP_ALIGN.CENTER
        _set_font(p_arr, 26, bold=True, color=C_ARROW)
        cx += ARRW

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — Общее сравнение моделей
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Общее сравнение моделей")

add_table(sl,
    ["Модель", "MAE", "RMSE", "Mean error"],
    [["CRR + GARCH",   "107.0", "145.8", "−78.1"],
     ["CRR + HV_63d",  "107.1", "140.6", "−89.7"],
     ["CRR + HV_21d",  "119.3", "156.4", "−95.3"]],
    left=Inches(0.4), top=Inches(1.05), width=Inches(12.5), height=Inches(1.9), font_size=17,
)

add_box(sl,
    "GARCH ≈ HV_63d по MAE  ·  HV_21d хуже  ·  Mean error < 0: все модели недооценивают рынок",
    left=Inches(0.4), top=Inches(3.05), width=Inches(12.5), height=Inches(0.5),
    fill=C_BOX_GREEN, font_size=14)

# fig_01_03: 2700x600 ratio=4.5 → width=12.5 → height=2.78
add_image(sl, FIGS/"fig_01_03_model_comparison.png",
          left=Inches(0.4), top=Inches(3.65), width=Inches(12.5), height=Inches(2.78))

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — TVNAE
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Почему стандартных метрик недостаточно: TVNAE")

# Left: formula + explanation
add_box(sl, "TVNAEᵢ = |P̂ᵢ − Pᵢ| / TVᵢ",
        left=Inches(0.4), top=Inches(1.05), width=Inches(6.0), height=Inches(0.6),
        fill=C_BOX_MAIN, font_size=20)
add_body(sl, ["где TVᵢ = Pᵢ − H(Fᵢ, Kᵢ)  (цена − внутренняя стоимость)"],
         top=Inches(1.73), width=Inches(6.0), font_size=13)
add_body(sl, [
    "MAE/RMSE считают ошибку по полной цене опциона.",
    "TVNAE нормирует на временну́ю стоимость — часть, которую модель объясняет.",
    "",
    "Deep ITM: 90% цены — внутренняя стоимость → MAE занижает масштаб ошибки.",
], top=Inches(2.1), width=Inches(6.0), font_size=15)

# Right: TVNAE table + conclusion
add_table(sl,
    ["Модель", "TVNAE mean", "TVNAE median"],
    [["CRR + GARCH",  "0.66", "0.62"],
     ["CRR + HV_63d", "0.70", "0.68"],
     ["CRR + HV_21d", "0.76", "0.78"]],
    left=Inches(6.7), top=Inches(1.05), width=Inches(6.3), height=Inches(1.7), font_size=16,
)
add_box(sl,
    "По TVNAE GARCH лучше описывает временну́ю стоимость\nDTE 0–7: GARCH TVNAE = 0.19 vs 0.59 у HV_63d",
    left=Inches(6.7), top=Inches(2.85), width=Inches(6.3), height=Inches(0.65),
    fill=C_BOX_GREEN, font_size=14)

# fig_pres_02_tvnae_heatmap: 1252x582 ratio=2.15 → height=3.4 → width=7.31
add_image(sl, FIGS/"fig_pres_02_tvnae_heatmap.png",
          left=Inches(3.01), top=Inches(3.6), width=Inches(7.31), height=Inches(3.4))

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — Американская vs Black-76
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Американская модель vs Black-76")

add_table(sl,
    ["Модель", "MAE", "RMSE", "Mean error"],
    [["CRR + HV_63d (амер.)",     "107.1", "140.6", "−89.7"],
     ["Black-76 + HV_63d (евр.)", "156.2", "211.6", "−149.0"]],
    left=Inches(0.4), top=Inches(1.05), width=Inches(12.5), height=Inches(1.5), font_size=17,
)

add_box(sl, "24.5% цен Black-76 ниже внутренней стоимости → прямой арбитраж",
        left=Inches(0.4), top=Inches(2.65), width=Inches(12.5), height=Inches(0.5),
        fill=C_BOX_RED, font_size=14)
add_box(sl, "Средняя премия раннего исполнения: 59.28 пункта  (медиана 9.87; high vol + DTE 91+: 106 пт)",
        left=Inches(0.4), top=Inches(3.24), width=Inches(12.5), height=Inches(0.5),
        fill=C_BOX_RED, font_size=14)
add_box(sl, "Класс модели важнее выбора волатильности",
        left=Inches(0.4), top=Inches(3.83), width=Inches(12.5), height=Inches(0.5),
        fill=C_BOX_GREEN, font_size=15)

# fig_pres_05: 1376x584 ratio=2.36 → height=3.0 → width=7.08, centered
add_image(sl, FIGS/"fig_pres_05_early_exercise_premium_heatmap.png",
          left=Inches(3.12), top=Inches(4.43), width=Inches(7.08), height=Inches(3.0))

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 9 — Ошибки по режимам
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Ошибки по рыночным режимам")

add_table(sl,
    ["Режим", "CRR + GARCH", "CRR + HV_63d", "CRR + HV_21d"],
    [["Low vol",  "77.5",  "82.1",  "104.1"],
     ["Mid vol",  "118.2", "128.3", "141.8"],
     ["High vol", "141.0", "122.8", "117.8"]],
    left=Inches(0.4), top=Inches(1.05), width=Inches(12.5), height=Inches(1.9), font_size=17,
)

add_body(sl, [
    "- Low / mid vol: GARCH лучший, HV_21d заметно хуже",
    "- High vol: преимущество GARCH исчезает — GARCH инерционен",
    "→ Режим влияет на ошибку описательно, но часть эффекта объясняется DTE и moneyness",
], top=Inches(3.05), width=Inches(12.5), font_size=16)

# fig_04_rolling_mae: 1500x600 ratio=2.5 → height=3.4 → width=8.5, centered
add_image(sl, FIGS/"fig_04_rolling_mae.png",
          left=Inches(2.42), top=Inches(3.85), width=Inches(8.5), height=Inches(3.4))

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 10 — Ошибки по сроку и денежности
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Ошибки по сроку до экспирации и денежности")

add_table(sl,
    ["DTE", "CRR + GARCH", "CRR + HV_63d", "CRR + HV_21d"],
    [["0–7 дней",   "1.88",  "3.53",  "2.78"],
     ["8–30 дней",  "4.87",  "9.60",  "8.79"],
     ["31–90 дней", "24.6",  "35.4",  "35.94"],
     ["91+ дней",   "122.5", "121.0", "135.3"]],
    left=Inches(0.4), top=Inches(1.05), width=Inches(12.5), height=Inches(2.1), font_size=17,
)

add_body(sl, [
    "- GARCH лучше на коротких и среднесрочных опционах",
    "- На длинных (91+): HV_63d устойчивее — GARCH теряет преимущество",
    "- Крупные ошибки deep ITM частично связаны с внутренней стоимостью — поэтому важен TVNAE",
], top=Inches(3.27), width=Inches(12.5), font_size=16)

# fig_01_03: 2700x600 ratio=4.5 → width=12.5 → height=2.78
add_image(sl, FIGS/"fig_01_03_model_comparison.png",
          left=Inches(0.4), top=Inches(4.05), width=Inches(12.5), height=Inches(2.78))

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 11 — Формальные проверки и robustness
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Формальные проверки и устойчивость")

# Left: statistical tests
add_box(sl, "Тест средней ошибки",
        left=Inches(0.4), top=Inches(1.05), width=Inches(6.1), height=Inches(0.42),
        fill=C_BOX_MAIN, font_size=13)
add_body(sl, [
    "Все три модели: t-stat от −20 до −25, p < 0.001",
    "→ Систематическое занижение статистически значимо",
], top=Inches(1.55), width=Inches(6.1), font_size=15)

add_box(sl, "Diebold–Mariano (попарно, HAC SE)",
        left=Inches(0.4), top=Inches(2.45), width=Inches(6.1), height=Inches(0.42),
        fill=C_BOX_MAIN, font_size=13)
add_body(sl, [
    "GARCH и HV_63d статистически неразличимы",
    "HV_21d значимо хуже HV_63d (p = 0.0003)",
], top=Inches(2.95), width=Inches(6.1), font_size=15)

add_box(sl, "OLS на абсолютных ошибках  (R² = 0.27)",
        left=Inches(0.4), top=Inches(4.15), width=Inches(6.1), height=Inches(0.42),
        fill=C_BOX_MAIN, font_size=13)
add_body(sl, [
    "DTE (+0.16, p<0.001) и |log-moneyness| (−178, p<0.001)",
    "объясняют бо́льшую часть вариации ошибок —",
    "больше, чем режим волатильности",
], top=Inches(4.65), width=Inches(6.1), font_size=15)

# Right: robustness table
add_box(sl, "Robustness",
        left=Inches(6.8), top=Inches(1.05), width=Inches(6.3), height=Inches(0.42),
        fill=C_BOX_ALT, font_size=13)
add_table(sl,
    ["Подвыборка", "Лучшая"],
    [["Полная выборка",       "GARCH ≈ HV_63d"],
     ["Без deep OTM / ITM",  "GARCH"],
     ["DTE ≤ 91",             "GARCH (22.27 vs 29.92)"],
     ["DTE > 91",             "HV_63d ≈ GARCH"],
     ["Без high_vol",         "GARCH"],
     ["r = 12% / 16% / 20%", "MAE меняется на ±5 пт"]],
    left=Inches(6.8), top=Inches(1.55), width=Inches(6.3), height=Inches(2.9), font_size=14,
)
add_body(sl, ["→ Результаты устойчивы во всех подвыборках"],
         top=Inches(4.65), left=Inches(6.8), width=Inches(6.3), font_size=15)

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 12 — Итоговые выводы
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Итоговые выводы")

conclusions = [
    ("1.", "Гипотеза подтверждается в уточнённой форме:\n"
           "режим влияет на ошибку, но DTE и moneyness объясняют значительную часть эффекта"),
    ("2.", "Для этих контрактов нужна американская модель:\n"
           "Black-76 даёт MAE на 46% хуже и 24.5% наблюдений ниже внутренней стоимости"),
    ("3.", "HV_63d — устойчивый и простой baseline:\n"
           "статистически неотличим от GARCH на полной выборке"),
    ("4.", "По TVNAE GARCH лучше на коротких DTE:\n"
           "TVNAE 0.66 vs 0.70; на DTE 0–7 ошибка в 3 раза меньше"),
    ("5.", "Все модели систематически недооценивают рынок:\n"
           "необходимы IV surface, liquidity adjustment, stochastic volatility, jumps"),
]

Y_c = Inches(1.1)
for num, text in conclusions:
    add_box(sl, num,
            left=Inches(0.4), top=Y_c, width=Inches(0.52), height=Inches(0.72),
            fill=C_BOX_MAIN, font_size=15)
    tb_c = sl.shapes.add_textbox(Inches(1.08), Y_c, W - Inches(1.5), Inches(0.72))
    tf_c = tb_c.text_frame; tf_c.word_wrap = True
    p_c  = tf_c.paragraphs[0]
    p_c.text = text
    _set_font(p_c, 15, color=C_TEXT)
    Y_c += Inches(1.02)

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 13 — Ограничения и направления развития
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()
add_title(sl, "Ограничения и направления развития")

add_box(sl, "Ограничения исследования",
        left=Inches(0.4), top=Inches(1.05), width=Inches(5.9), height=Inches(0.42),
        fill=C_BOX_MAIN, font_size=13)
add_body(sl, [
    "- Расчётная цена биржи, а не котировка — на неликвидных",
    "  инструментах эти значения могут расходиться",
    "- Выборка смещена к длинным опционам в деньгах —",
    "  ограничивает перенос выводов на другие рынки",
    "- Единая σ без поверхности — систематическое смещение",
    "  неустранимо в данной постановке",
    "- Фиксированная ставка не учитывает срочную структуру",
], top=Inches(1.55), width=Inches(5.9), font_size=15)

add_box(sl, "Направления развития",
        left=Inches(7.0), top=Inches(1.05), width=Inches(6.0), height=Inches(0.42),
        fill=C_BOX_ALT, font_size=13)
add_body(sl, [
    "- Переход к поверхности подразумеваемой волатильности (IV surface)",
    "- Учёт ликвидностных факторов в модели",
    "- Модели со стохастической волатильностью (Heston, SABR)",
    "- Модели со скачками (jumps) — мотивированы тяжёлыми хвостами",
    "- Срочная структура ставок вместо единого значения",
], top=Inches(1.55), left=Inches(7.0), width=Inches(6.0), font_size=15)

# ═══════════════════════════════════════════════════════════════════════════
# СЛАЙД 14 — Спасибо
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()

bar2 = sl.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                            Inches(0), Inches(0), W, Inches(0.18))
bar2.fill.solid(); bar2.fill.fore_color.rgb = C_BOX_MAIN
bar2.line.fill.background()

tb = sl.shapes.add_textbox(Inches(1), Inches(2.4), W - Inches(2), Inches(1.2))
p  = tb.text_frame.paragraphs[0]
p.text = "Спасибо за внимание!"
_set_font(p, 40, bold=True, color=C_TITLE)
p.alignment = PP_ALIGN.CENTER

tb2 = sl.shapes.add_textbox(Inches(1), Inches(3.8), W - Inches(2), Inches(0.6))
p2  = tb2.text_frame.paragraphs[0]
p2.text = "Вопросы?"
_set_font(p2, 24, color=RGBColor(0x55, 0x55, 0x55))
p2.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════════
out = FINAL / "presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({out.stat().st_size // 1024} KB)")
